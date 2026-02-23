"""
Generate a presentation deck for the Deep Analysis – Code Analyser platform.
Content is derived from the actual codebase and architecture documentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x11, 0x1A)   # slide background
ACCENT     = RGBColor(0x64, 0x6C, 0xFF)   # indigo accent
ACCENT2    = RGBColor(0x00, 0xD9, 0xFF)   # cyan accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
MUTED      = RGBColor(0x88, 0x88, 0xAA)
RED_ACCENT = RGBColor(0xFF, 0x4D, 0x6A)
GREEN_ACC  = RGBColor(0x00, 0xE6, 0x96)
ORANGE_ACC = RGBColor(0xFF, 0xA7, 0x26)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ARCH_IMG   = os.path.join(SCRIPT_DIR, "svgviewer-png-output.png")
OUTPUT     = os.path.join(SCRIPT_DIR, "Deep_Analysis_Deck.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, opacity=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Rounded corners
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_frame(slide, left, top, width, height, items,
                     font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_items,
             card_color=RGBColor(0x1A, 0x1D, 0x2E), title_color=ACCENT2):
    card = add_shape_rect(slide, left, top, width, height, card_color)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.15),
                 width - Inches(0.5), Inches(0.4),
                 title, font_size=15, bold=True, color=title_color)
    add_bullet_frame(slide, left + Inches(0.25), top + Inches(0.55),
                     width - Inches(0.5), height - Inches(0.7),
                     body_items, font_size=13, color=LIGHT_GRAY)
    return card


# ── Slide builders ──────────────────────────────────────────────────
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Accent line at top
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, ACCENT)

    # Title
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                 "Deep Analysis", font_size=54, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.8),
                 "Code Analyser Platform", font_size=36, bold=False,
                 color=ACCENT2, alignment=PP_ALIGN.CENTER)

    # Tagline
    add_text_box(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.7),
                 "Understand HOW you solve problems — not just whether you succeed.",
                 font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Subtitle chips
    chips = ["AST-Powered Analysis", "Real-Time Code Execution", "Personalized Insights"]
    chip_w = Inches(3)
    start_x = Inches(2)
    for i, chip in enumerate(chips):
        x = start_x + i * (chip_w + Inches(0.5))
        add_shape_rect(slide, x, Inches(5.3), chip_w, Inches(0.55),
                       RGBColor(0x1A, 0x1D, 0x2E))
        add_text_box(slide, x, Inches(5.33), chip_w, Inches(0.5),
                     chip, font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # Bottom line
    add_accent_line(slide, Inches(0), Inches(7.2), SLIDE_W, ACCENT)


def slide_problem_statement(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "The Problem", font_size=36, bold=True, color=RED_ACCENT)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), RED_ACCENT)

    # Main statement
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8),
                 "Existing coding platforms (LeetCode, HackerRank, CodeChef) operate on binary feedback: Pass or Fail.",
                 font_size=22, color=WHITE)

    # What they don't know – cards
    gaps = [
        ("❌  Syntax vs. Logic?",   "Platforms don't distinguish whether a failure is a typo\nor a fundamental misunderstanding of the algorithm."),
        ("❌  Concepts Understood?", "No visibility into which data structures or patterns\nthe student actually grasped during their attempts."),
        ("❌  What Changed?",        "No tracking of how the code evolved across\nmultiple attempts — the iteration story is lost."),
        ("❌  What to Practice?",    "No personalised next-step recommendations;\nstudents are left guessing what to study next."),
    ]
    card_w = Inches(2.75)
    card_h = Inches(2.0)
    start_x = Inches(0.8)
    y = Inches(2.5)
    for i, (title, body) in enumerate(gaps):
        x = start_x + i * (card_w + Inches(0.2))
        add_card(slide, x, y, card_w, card_h, title,
                 body.split("\n"), card_color=RGBColor(0x25, 0x15, 0x1A),
                 title_color=RED_ACCENT)

    # Impact statement
    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.6),
                 "Impact:  Students repeat the same mistakes, instructors lack diagnostic data,\n"
                 "and bootcamps cannot measure real learning outcomes.\n\n"
                 "The core issue is that coding platforms were designed for assessment — not for learning.",
                 font_size=17, color=MUTED)


def slide_our_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Our Solution", font_size=36, bold=True, color=GREEN_ACC)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), GREEN_ACC)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.7),
                 "Deep behavioural analysis of every coding session — tracking, analysing, and generating actionable insights.",
                 font_size=20, color=WHITE)

    # Three pillars
    pillars = [
        ("1 ▸ Track Everything",
         ["Every code snapshot saved immutably",
          "Time-stamped attempt timeline",
          "Error progression across submissions",
          "Test case pass/fail patterns",
          "Multi-language support (Python, C, C++, Java)"]),
        ("2 ▸ Analyse Deeply",
         ["AST parsing detects concepts used",
          "Code diffs between consecutive attempts",
          "Error classification (syntax / logic / runtime)",
          "Time-complexity heuristic estimation",
          "Topic mastery confidence scoring",
          "Edit-hotspot detection (line-level)"]),
        ("3 ▸ Generate Insights",
         ["Natural-language summary of the session",
          "Concept-level breakdown with confidence %",
          "Personalised practice recommendations",
          "Weak-topic cards (Easy → Medium → Hard)",
          "Monthly SWOT report per student",
          "Dashboard KPIs & progress tracking"]),
    ]
    card_w = Inches(3.7)
    card_h = Inches(3.5)
    start_x = Inches(0.8)
    y = Inches(2.3)
    colors = [ACCENT, ACCENT2, GREEN_ACC]
    for i, (title, items) in enumerate(pillars):
        x = start_x + i * (card_w + Inches(0.2))
        add_card(slide, x, y, card_w, card_h, title, items,
                 card_color=RGBColor(0x14, 0x1D, 0x1A), title_color=colors[i])

    add_text_box(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.5),
                 "Result:  \"You understood loops and conditionals but struggled with edge-case handling. Practice: Two Sum (Easy) → 3Sum (Medium).\"",
                 font_size=15, color=MUTED)


def slide_how_it_works(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "How It Works — User Flow", font_size=36, bold=True, color=ACCENT)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), ACCENT)

    steps = [
        ("① Login",         "Firebase Auth\nEmail / Google SSO",           ACCENT),
        ("② Browse",        "Filter by difficulty,\ntag, or search",       ACCENT2),
        ("③ Write Code",    "Monaco Editor\nAuto-save every 30s",          GREEN_ACC),
        ("④ Run / Submit",  "Judge0 sandbox\nexecution",                   ORANGE_ACC),
        ("⑤ Analyse",       "AST + Diff + Error\nclassification",          ACCENT),
        ("⑥ View Insights", "Summary, concepts,\nrecommendations",         ACCENT2),
        ("⑦ Dashboard",     "KPIs, SWOT,\nmonthly reports",                GREEN_ACC),
    ]
    card_w = Inches(1.55)
    card_h = Inches(2.4)
    start_x = Inches(0.5)
    y = Inches(1.6)
    for i, (title, body, col) in enumerate(steps):
        x = start_x + i * (card_w + Inches(0.15))
        c = add_shape_rect(slide, x, y, card_w, card_h, RGBColor(0x1A, 0x1D, 0x2E))
        add_text_box(slide, x, y + Inches(0.2), card_w, Inches(0.45),
                     title, font_size=15, bold=True, color=col,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.75), card_w - Inches(0.2), Inches(1.5),
                     body, font_size=12, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # arrow connectors (simple rectangles)
    for i in range(len(steps) - 1):
        x = start_x + (i + 1) * (card_w + Inches(0.15)) - Inches(0.12)
        add_shape_rect(slide, x, y + Inches(1.0), Inches(0.09), Inches(0.35),
                       RGBColor(0x44, 0x44, 0x66))

    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5),
                 "Key implementation details from the codebase:\n\n"
                 "• ProblemWorkspacePage.tsx (1 518 lines) — full workspace with Monaco editor, "
                 "language selector, run/submit flows, inline test results, and side-by-side code diff viewer.\n"
                 "• DashboardPage.tsx — KPIs (solved, attempted, success rate), concept mastery bars, "
                 "recent sessions deep-link, and monthly SWOT report generation.\n"
                 "• Analysis engine returns concept breakdown with confidence %, attempt timeline, "
                 "weak-topic diagnostics, and targeted practice suggestions.\n"
                 "• Multi-language support: Python 3.x, C (GCC), C++ (GCC), Java (OpenJDK).",
                 font_size=13, color=MUTED)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "System Architecture", font_size=36, bold=True, color=ACCENT)
    add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(2.5), ACCENT)

    # Architecture image
    if os.path.exists(ARCH_IMG):
        slide.shapes.add_picture(ARCH_IMG,
                                 Inches(0.5), Inches(1.15),
                                 width=Inches(8.5))

    # Tech stack sidebar
    stack_items = [
        ("Frontend", "React 18 · TypeScript · Vite\nMonaco Editor · Zustand\nFirebase Auth", ACCENT),
        ("Backend", "Node.js 20 · Express 4\nPrisma ORM · JWT\nAzure Container Apps", ACCENT2),
        ("Analysis", "Python 3.11 · FastAPI\nAST Parser · difflib · radon\nAzure Container Apps", GREEN_ACC),
        ("Database", "PostgreSQL 15\nNeon (serverless)\nPrisma migrations", ORANGE_ACC),
        ("Execution", "Judge0 CE\nDocker · isolate sandbox\nAzure VM (self-hosted)", RED_ACCENT),
    ]
    card_w = Inches(3.5)
    card_h = Inches(0.95)
    sx = Inches(9.4)
    for i, (title, body, col) in enumerate(stack_items):
        y = Inches(1.15) + i * (card_h + Inches(0.12))
        c = add_shape_rect(slide, sx, y, card_w, card_h, RGBColor(0x1A, 0x1D, 0x2E))
        add_text_box(slide, sx + Inches(0.15), y + Inches(0.05),
                     card_w - Inches(0.3), Inches(0.3),
                     title, font_size=13, bold=True, color=col)
        add_text_box(slide, sx + Inches(0.15), y + Inches(0.35),
                     card_w - Inches(0.3), Inches(0.55),
                     body, font_size=10, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
                 "Zero-cost hosting  •  All services on free tiers  •  Horizontally scalable micro-service design",
                 font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)


def slide_tech_deep_dive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Analysis Engine — Deep Dive", font_size=36, bold=True, color=ACCENT2)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), ACCENT2)

    left_items = [
        ("AST Concept Detection",
         ["Parses Python/C/C++/Java code into abstract syntax trees",
          "Detects: loops, conditionals, recursion, data structures",
          "Assigns confidence scores per concept"]),
        ("Code Diff Engine",
         ["Side-by-side diff with line-level change tracking",
          "buildDiffRows() generates same / added / removed / changed rows",
          "Visualises iteration strategy across attempts"]),
        ("Error Classification",
         ["Maps Judge0 status codes to error categories",
          "Tracks syntax / logic / runtime error rates across attempts",
          "Identifies dominant error pattern per session"]),
    ]
    right_items = [
        ("Time Complexity Heuristic",
         ["Rule-based estimation from AST loop/recursion depth",
          "radon cyclomatic complexity scoring",
          "Suggests optimisation opportunities"]),
        ("Insight Generation",
         ["Natural language summary template engine",
          "Concept breakdown table with confidence %",
          "Weak-topic extraction with practice-next cards",
          "Edit-hotspot detection at line level"]),
        ("Monthly SWOT Reports",
         ["Strengths / Weaknesses / Opportunities / Threats",
          "Concept mastery progression over time",
          "Aggregated across all sessions in period"]),
    ]

    card_w = Inches(5.7)
    card_h = Inches(1.65)
    for i, (title, items) in enumerate(left_items):
        y = Inches(1.4) + i * (card_h + Inches(0.12))
        add_card(slide, Inches(0.5), y, card_w, card_h, title, items,
                 title_color=ACCENT2)
    for i, (title, items) in enumerate(right_items):
        y = Inches(1.4) + i * (card_h + Inches(0.12))
        add_card(slide, Inches(6.5), y, card_w, card_h, title, items,
                 title_color=GREEN_ACC)


def slide_competitive_advantage(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Competitive Advantage", font_size=36, bold=True, color=ORANGE_ACC)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), ORANGE_ACC)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
                 "Why can't LeetCode / HackerRank simply copy this?",
                 font_size=22, color=WHITE)

    advantages = [
        ("Architectural Lock-in",
         "Existing platforms are built around single-submission evaluation.\n"
         "Retrofitting multi-attempt tracking requires fundamental re-architecture.",
         RED_ACCENT),
        ("Missing Data",
         "Competitors discard intermediate attempts — they only store final\n"
         "submissions. We capture every snapshot, creating a rich behavioral dataset.",
         ACCENT),
        ("Different Business Model",
         "Their revenue comes from hiring assessments, not education.\n"
         "Our model is education-first: freemium for individuals, B2B for bootcamps.",
         GREEN_ACC),
        ("First-Mover Data Moat",
         "Every session builds our behavioral dataset. Over time this compounds\n"
         "into ML-ready training data for predictive learning recommendations.",
         ORANGE_ACC),
    ]
    card_w = Inches(5.7)
    card_h = Inches(1.35)
    for i, (title, body, col) in enumerate(advantages):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(0.8) + col_idx * (card_w + Inches(0.4))
        y = Inches(2.2) + row_idx * (card_h + Inches(0.25))
        add_card(slide, x, y, card_w, card_h, title,
                 body.split("\n"), title_color=col)


def slide_demo_and_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Demo Flow & Future Roadmap", font_size=36, bold=True, color=ACCENT)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), ACCENT)

    # Demo flow
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
                 "Live Demo (5 min)", font_size=22, bold=True, color=ACCENT2)
    demo_steps = [
        "1.  Login as demo user (Firebase Auth)",
        "2.  Browse problem catalog → select \"Two Sum\"",
        "3.  Write code in Monaco editor (Python / C++ / Java)",
        "4.  Run against visible tests → see inline results",
        "5.  Submit 3 attempts with different errors",
        "6.  Click \"Analyse\" → reveal deep insights panel",
        "7.  Navigate to Dashboard → KPIs + SWOT report",
    ]
    add_bullet_frame(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5),
                     demo_steps, font_size=15, color=LIGHT_GRAY)

    # Future roadmap
    add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.4),
                 "Roadmap — Phase 2+", font_size=22, bold=True, color=GREEN_ACC)
    roadmap = [
        "◆  ML-powered recommendation engine",
        "◆  Cross-problem misconception detection",
        "◆  Adaptive difficulty learning paths",
        "◆  JavaScript / Go / Rust support",
        "◆  Expand to 70+ curated problems",
        "◆  Instructor analytics dashboard",
        "◆  OAuth (GitHub, Microsoft) login",
        "◆  Real-time collaborative code review",
    ]
    add_bullet_frame(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(3.5),
                     roadmap, font_size=15, color=LIGHT_GRAY)

    # Bottom bar
    add_shape_rect(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.3),
                   RGBColor(0x1A, 0x1D, 0x2E))
    add_text_box(slide, Inches(0.8), Inches(5.85), Inches(11.8), Inches(0.4),
                 "Business Model", font_size=18, bold=True, color=ORANGE_ACC)
    add_text_box(slide, Inches(0.8), Inches(6.25), Inches(11.8), Inches(0.6),
                 "Freemium:  Free basic access  ·  Premium for ML insights & advanced reports  |  "
                 "B2B:  Bootcamp & university licenses with instructor dashboards",
                 font_size=14, color=LIGHT_GRAY)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, ACCENT)

    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
                 "Deep Analysis", font_size=48, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
                 "Understand your code. Improve your craft.",
                 font_size=26, color=ACCENT2, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(4.5), Inches(9), Inches(0.6),
                 "\"Existing platforms tell you PASS or FAIL.\nWe show you HOW you solve problems — and what to do next.\"",
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.5),
                 "Thank you!",
                 font_size=32, bold=True, color=GREEN_ACC, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(0), Inches(7.2), SLIDE_W, ACCENT)


# ── Main ────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)              # 1: Title
    slide_problem_statement(prs)  # 2: Problem
    slide_our_solution(prs)       # 3: Solution
    slide_how_it_works(prs)       # 4: User flow
    slide_architecture(prs)       # 5: Architecture + image
    slide_tech_deep_dive(prs)     # 6: Analysis engine deep dive
    slide_competitive_advantage(prs)  # 7: Why us
    slide_demo_and_future(prs)    # 8: Demo + Roadmap
    slide_closing(prs)            # 9: Closing

    prs.save(OUTPUT)
    print(f"✅ Presentation saved to: {OUTPUT}")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
